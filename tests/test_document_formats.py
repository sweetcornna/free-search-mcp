"""read_doc format coverage beyond pdf/docx/html/text.

Everything here is generated in-process rather than checked in as a binary
fixture, so the tests stay readable and can't rot into "some .xlsx from 2024".
"""

from __future__ import annotations

import io
import tarfile
import zipfile
from pathlib import Path

import pytest

from search_mcp import config, documents
from search_mcp.documents import (
    _detect_format,
    _parse_archive,
    _parse_code,
    _parse_csv,
    read_document,
)

# pytest.ini sets `asyncio_mode = auto` so async tests are auto-marked.


@pytest.fixture
def sandbox(tmp_path, monkeypatch):
    """Point the local-read sandbox at tmp_path (reads are off by default)."""
    monkeypatch.setattr(config.settings, "document_root", tmp_path)
    monkeypatch.setattr(documents.settings, "document_root", tmp_path)
    return tmp_path


# ---------------------------------------------------------------------------
# Format detection
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "source,expected",
    [
        ("a.pdf", "pdf"), ("a.docx", "docx"), ("a.xlsx", "xlsx"),
        ("a.pptx", "pptx"), ("a.epub", "epub"), ("a.csv", "csv"),
        ("a.tsv", "csv"), ("a.zip", "archive"), ("a.tar.gz", "archive"),
        ("a.whl", "archive"), ("a.png", "image"), ("a.svg", "image"),
        ("a.heic", "image"), ("a.py", "code"), ("a.rs", "code"),
        ("a.json", "code"), ("a.yaml", "code"), ("Dockerfile", "code"),
        ("Makefile", "code"), ("a.html", "html"), ("a.md", "markdown"),
        ("a.txt", "text"), ("a.bin", "unknown"),
    ],
)
def test_detect_format_by_extension(source, expected):
    assert _detect_format(source) == expected


@pytest.mark.parametrize(
    "ctype,expected",
    [
        ("application/pdf", "pdf"),
        ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx"),
        ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"),
        ("application/epub+zip", "epub"),
        ("text/csv", "csv"),
        ("application/zip", "archive"),
        ("image/png", "image"),
        ("image/svg+xml", "image"),
        ("text/plain; charset=utf-8", "text"),
    ],
)
def test_detect_format_by_content_type(ctype, expected):
    """An extension-less URL still has to be classified from its headers."""
    assert _detect_format("https://x.example/download?id=1", ctype) == expected


def test_query_string_does_not_confuse_extension_detection():
    assert _detect_format("https://x.example/data.csv?token=abc.png") == "csv"


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------


def test_csv_becomes_a_markdown_table():
    text, truncated = _parse_csv(b"name,qty\nwidget,3\n")
    assert text.splitlines()[0] == "| name | qty |"
    assert text.splitlines()[1] == "| --- | --- |"
    assert truncated is False


def test_csv_escapes_pipes_so_cells_cannot_forge_columns():
    text, _ = _parse_csv(b"name,qty\ngad|get,4\n")
    assert r"gad\|get" in text


def test_csv_detects_tab_delimiters():
    text, _ = _parse_csv(b"name\tqty\nwidget\t3\n")
    assert "| name | qty |" in text


def test_csv_ragged_rows_are_padded():
    """A short row must not shift the remaining columns."""
    text, _ = _parse_csv(b"a,b,c\n1\n")
    assert "| 1 |  |  |" in text


def test_csv_is_capped_by_max_document_chars(monkeypatch):
    monkeypatch.setattr(documents.settings, "max_document_chars", 20)
    body = b"a,b\n" + b"xxxx,yyyy\n" * 500
    _text, truncated = _parse_csv(body)
    assert truncated is True


# ---------------------------------------------------------------------------
# Code
# ---------------------------------------------------------------------------


def test_code_is_fenced_with_its_language():
    out = _parse_code(b"print(1)\n", "x.py")
    assert out.startswith("```python\n")


def test_code_fence_widens_past_backticks_in_the_file():
    """A file containing a Markdown fence must not break out of ours."""
    out = _parse_code(b"a = '''```'''\n", "x.py")
    assert out.startswith("````python")
    assert out.rstrip().endswith("````")


def test_code_unknown_extension_still_fences_without_a_language():
    out = _parse_code(b"data\n", "x.unknownext")
    assert out.startswith("```\n")


# ---------------------------------------------------------------------------
# Archives — listed, never extracted
# ---------------------------------------------------------------------------


def _zip_bytes(entries: dict[str, bytes]) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data in entries.items():
            z.writestr(name, data)
    return buf.getvalue()


def test_zip_is_listed_not_extracted():
    blob = _zip_bytes({"a.txt": b"x" * 100, "dir/b.txt": b"y" * 50})
    text, _ = _parse_archive(blob, "a.zip")
    assert "a.txt" in text and "dir/b.txt" in text
    assert "listed, not extracted" in text
    # The payload itself must not appear in the output.
    assert "xxxxx" not in text


def test_tar_is_listed_too():
    buf = io.BytesIO()
    with tarfile.open(fileobj=buf, mode="w") as tf:
        info = tarfile.TarInfo("inside.txt")
        data = b"z" * 40
        info.size = len(data)
        tf.addfile(info, io.BytesIO(data))
    text, _ = _parse_archive(buf.getvalue(), "a.tar")
    assert "inside.txt" in text


def test_zip_bomb_ratio_is_called_out():
    """A highly compressible archive is exactly the zip-bomb shape; the
    listing has to say so rather than looking like an ordinary archive."""
    blob = _zip_bytes({"bomb.txt": b"\0" * 2_000_000})
    text, _ = _parse_archive(blob, "bomb.zip")
    assert "expands more than 100x" in text


def test_archive_entry_count_is_capped():
    blob = _zip_bytes({f"f{i}.txt": b"x" for i in range(2100)})
    _text, truncated = _parse_archive(blob, "many.zip")
    assert truncated is True


def test_corrupt_archive_reports_instead_of_raising():
    text, truncated = _parse_archive(b"not an archive at all", "broken.zip")
    assert "Could not read archive" in text
    assert truncated is False


# ---------------------------------------------------------------------------
# Office formats, end to end through read_document
# ---------------------------------------------------------------------------


async def test_xlsx_renders_one_table_per_sheet(sandbox: Path):
    from openpyxl import Workbook

    wb = Workbook()
    first = wb.active
    first.title = "Alpha"
    first.append(["h1", "h2"])
    first.append(["a", 1])
    second = wb.create_sheet("Beta")
    second.append(["only"])
    wb.save(sandbox / "b.xlsx")

    result = await read_document(str(sandbox / "b.xlsx"))
    assert result.format == "xlsx"
    assert "## Alpha" in result.content
    assert "## Beta" in result.content
    assert "| h1 | h2 |" in result.content


async def test_xlsx_skips_the_empty_padding_rows_spreadsheets_carry(sandbox: Path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["h"])
    ws.append([None])
    ws.append(["v"])
    wb.save(sandbox / "pad.xlsx")

    result = await read_document(str(sandbox / "pad.xlsx"))
    # header + separator + the one real data row; the blank row is dropped.
    rows = [ln for ln in result.content.splitlines() if ln.startswith("|")]
    assert rows == ["| h |", "| --- |", "| v |"]


async def test_pptx_reports_slide_count_and_notes(sandbox: Path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Title here"
    slide.placeholders[1].text = "Body text"
    slide.notes_slide.notes_text_frame.text = "remember this"
    prs.save(sandbox / "d.pptx")

    result = await read_document(str(sandbox / "d.pptx"))
    assert result.format == "pptx"
    assert result.pages == 1
    assert "## Slide 1" in result.content
    assert "Title here" in result.content
    assert "remember this" in result.content


async def test_code_file_round_trips(sandbox: Path):
    (sandbox / "m.py").write_text("def hi():\n    return 1\n", encoding="utf-8")
    result = await read_document(str(sandbox / "m.py"))
    assert result.format == "code"
    assert "```python" in result.content


async def test_csv_file_round_trips(sandbox: Path):
    (sandbox / "t.csv").write_text("name,qty\nwidget,3\n", encoding="utf-8")
    result = await read_document(str(sandbox / "t.csv"))
    assert result.format == "csv"
    assert "| name | qty |" in result.content


async def test_zip_file_round_trips(sandbox: Path):
    (sandbox / "a.zip").write_bytes(_zip_bytes({"x.txt": b"hello"}))
    result = await read_document(str(sandbox / "a.zip"))
    assert result.format == "archive"
    assert "x.txt" in result.content


async def test_images_are_redirected_to_fetch(sandbox: Path):
    """read_doc extracts text; an image has none. Returning an empty document
    would look like a parse failure."""
    (sandbox / "p.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"\0" * 32)
    with pytest.raises(ValueError, match="Use `fetch`"):
        await read_document(str(sandbox / "p.png"))


async def test_unsupported_format_lists_what_is_supported(sandbox: Path):
    (sandbox / "x.bin").write_bytes(b"\x00\x01\x02")
    with pytest.raises(ValueError, match="xlsx"):
        await read_document(str(sandbox / "x.bin"))


async def test_office_container_serving_html_falls_back(sandbox: Path):
    """A paywall or login wall served at a .xlsx URL must surface its text
    rather than crashing the parser."""
    (sandbox / "wall.xlsx").write_bytes(b"<html><body><h1>Please log in</h1></body></html>")
    result = await read_document(str(sandbox / "wall.xlsx"))
    assert result.format == "html"
    assert "Please log in" in result.content
