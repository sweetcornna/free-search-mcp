from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import httpx
from docx import Document as DocxDocument
from markdownify import markdownify as html_to_md
from pypdf import PdfReader

from .config import settings
from .formatting import estimate_tokens, smart_truncate
from .httpfetch import httpx_client_kwargs, httpx_stream_capped
from .url_safety import assert_url_allowed_async

log = logging.getLogger(__name__)


@dataclass(slots=True)
class DocumentResult:
    source: str
    format: str
    title: str
    content: str
    truncated: bool
    pages: int | None = None
    tokens_estimated: int = 0
    total_chars: int = 0
    start: int = 0
    returned_chars: int = 0

    def to_dict(self) -> dict[str, Any]:
        return {
            "source": self.source,
            "format": self.format,
            "title": self.title,
            "content": self.content,
            "truncated": self.truncated,
            "pages": self.pages,
            "tokens_estimated": self.tokens_estimated,
            "total_chars": self.total_chars,
            "start": self.start,
            "returned_chars": self.returned_chars,
        }


# Source-code and config extensions read as plain text, but tagged with the
# language so the renderer can fence them correctly. Mapping to the identifier
# a Markdown renderer expects, which is not always the extension.
_CODE_LANGS = {
    "py": "python", "pyi": "python", "js": "javascript", "mjs": "javascript",
    "cjs": "javascript", "ts": "typescript", "tsx": "tsx", "jsx": "jsx",
    "rb": "ruby", "go": "go", "rs": "rust", "java": "java", "kt": "kotlin",
    "swift": "swift", "c": "c", "h": "c", "cc": "cpp", "cpp": "cpp",
    "hpp": "cpp", "cs": "csharp", "php": "php", "pl": "perl", "lua": "lua",
    "r": "r", "scala": "scala", "sh": "bash", "bash": "bash", "zsh": "bash",
    "fish": "fish", "ps1": "powershell", "sql": "sql", "graphql": "graphql",
    "proto": "protobuf", "tf": "hcl", "hcl": "hcl", "dockerfile": "dockerfile",
    "makefile": "makefile", "cmake": "cmake", "gradle": "groovy",
    "json": "json", "jsonl": "json", "yaml": "yaml", "yml": "yaml",
    "toml": "toml", "ini": "ini", "cfg": "ini", "conf": "ini", "env": "bash",
    "xml": "xml", "svg": "xml", "css": "css", "scss": "scss", "less": "less",
    "vue": "vue", "svelte": "svelte", "tex": "latex", "diff": "diff",
    "patch": "diff",
}

_ARCHIVE_EXTS = (".zip", ".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tar.xz", ".whl", ".jar")

_IMAGE_EXTS = (
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif",
    ".ico", ".avif", ".heic", ".svg",
)


def _ext(source: str) -> str:
    """Lower-cased extension without the dot, ignoring any query string."""
    path = urlparse(source).path if "://" in source else source
    name = path.rsplit("/", 1)[-1].lower()
    # Extension-less build files still have a well-known identity.
    if name in ("dockerfile", "makefile", "cmakelists.txt"):
        return name.split(".", 1)[0]
    _, _, ext = name.rpartition(".")
    return ext if ext != name else ""


def _detect_format(source: str, content_type: str | None = None) -> str:
    # Match on the PATH, not the raw source: a query string routinely carries
    # something that looks like an extension ("...data.csv?token=abc.png"),
    # and endswith() on the whole URL would classify by the wrong one.
    s = (urlparse(source).path if "://" in source else source).lower()
    ctype = (content_type or "").lower()
    if s.endswith(".pdf") or "pdf" in ctype:
        return "pdf"
    if s.endswith(".docx") or "wordprocessingml" in ctype:
        return "docx"
    if s.endswith(".xlsx") or s.endswith(".xlsm") or "spreadsheetml" in ctype:
        return "xlsx"
    if s.endswith(".pptx") or "presentationml" in ctype:
        return "pptx"
    if s.endswith(".epub") or "epub" in ctype:
        return "epub"
    if s.endswith(".csv") or s.endswith(".tsv") or "text/csv" in ctype:
        return "csv"
    if s.endswith(_ARCHIVE_EXTS) or "zip" in ctype or "x-tar" in ctype:
        return "archive"
    # Images are described, not transcribed — see fetch's asset path.
    if s.endswith(_IMAGE_EXTS) or ctype.startswith("image/"):
        return "image"
    if s.endswith((".html", ".htm")) or "html" in ctype:
        return "html"
    if s.endswith((".md", ".markdown")):
        return "markdown"
    if _ext(source) in _CODE_LANGS:
        return "code"
    if s.endswith((".txt", ".log")) or ctype.startswith("text/"):
        return "text"
    return "unknown"


def _parse_pdf(blob: bytes) -> tuple[str, str, int, bool]:
    """Parse a PDF, capping pages and total text to defuse decompression bombs.

    Stops after ``settings.max_pdf_pages`` pages OR once accumulated text passes
    ``settings.max_document_chars``. Returns
    ``(title, text, total_pages, truncated)`` where ``truncated`` is True when
    either cap was hit (so callers can flag the result as incomplete).
    """
    reader = PdfReader(io.BytesIO(blob))
    total_pages = len(reader.pages)
    max_pages = settings.max_pdf_pages
    max_chars = settings.max_document_chars
    parts: list[str] = []
    acc_chars = 0
    truncated = False
    for i, page in enumerate(reader.pages, 1):
        if i > max_pages:
            truncated = True
            break
        try:
            txt = page.extract_text() or ""
        except Exception as e:
            log.warning("pdf page %d failed: %s", i, e)
            continue
        if txt.strip():
            piece = f"## Page {i}\n\n{txt.strip()}"
            parts.append(piece)
            acc_chars += len(piece)
            if acc_chars >= max_chars:
                truncated = True
                break
    title = ""
    try:
        meta = reader.metadata
        if meta and meta.title:
            title = str(meta.title)
    except Exception:
        pass
    return title, "\n\n".join(parts), total_pages, truncated


def _parse_docx(blob: bytes) -> tuple[str, bool]:
    """Parse a docx, capping accumulated text to defuse decompression bombs.

    Stops once accumulated text passes ``settings.max_document_chars``. Returns
    ``(text, truncated)``.
    """
    doc = DocxDocument(io.BytesIO(blob))
    max_chars = settings.max_document_chars
    parts: list[str] = []
    acc_chars = 0
    truncated = False
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = (p.style.name if p.style else "") or ""
        if style.startswith("Heading"):
            level = "".join(c for c in style if c.isdigit()) or "1"
            piece = f"{'#' * int(level)} {text}"
        else:
            piece = text
        parts.append(piece)
        acc_chars += len(piece)
        if acc_chars >= max_chars:
            truncated = True
            break
    if not truncated:
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                piece = "\n".join(rows)
                parts.append(piece)
                acc_chars += len(piece)
                if acc_chars >= max_chars:
                    truncated = True
                    break
    return "\n\n".join(parts), truncated


def _parse_html(blob: bytes) -> str:
    html = blob.decode("utf-8", errors="replace")
    return html_to_md(html, heading_style="ATX", bullets="-").strip()


def _parse_text(blob: bytes) -> str:
    return blob.decode("utf-8", errors="replace")


def _parse_code(blob: bytes, source: str) -> str:
    """Source/config file, fenced with its language so it renders as code."""
    lang = _CODE_LANGS.get(_ext(source), "")
    text = blob.decode("utf-8", errors="replace")
    # Pick a fence longer than any run of backticks inside the file, or a file
    # containing a Markdown fence would break out of ours.
    longest = 0
    run = 0
    for ch in text:
        run = run + 1 if ch == "`" else 0
        longest = max(longest, run)
    fence = "`" * max(3, longest + 1)
    return f"{fence}{lang}\n{text}\n{fence}"


def _rows_to_markdown(rows: list[list[str]], max_cols: int = 30) -> str:
    """Render rows as a Markdown table, treating row 0 as the header."""
    if not rows:
        return ""
    width = min(max(len(r) for r in rows), max_cols)
    if width == 0:
        return ""

    def line(cells: list[str]) -> str:
        padded = list(cells[:width]) + [""] * (width - len(cells[:width]))
        # Escape pipes so a cell can't fabricate extra columns.
        return "| " + " | ".join(c.replace("|", "\\|") for c in padded) + " |"

    out = [line(rows[0]), "|" + "|".join([" --- "] * width) + "|"]
    out.extend(line(r) for r in rows[1:])
    return "\n".join(out)


def _parse_csv(blob: bytes) -> tuple[str, bool]:
    """CSV/TSV to a Markdown table, capped at max_document_chars."""
    import csv as _csv

    text = blob.decode("utf-8", errors="replace")
    sample = text[:8192]
    try:
        dialect = _csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except _csv.Error:
        # A single-column file gives the sniffer nothing to work with.
        dialect = _csv.excel
    rows: list[list[str]] = []
    acc = 0
    truncated = False
    for row in _csv.reader(io.StringIO(text), dialect):
        cells = [str(c).strip() for c in row]
        acc += sum(len(c) for c in cells)
        rows.append(cells)
        if acc >= settings.max_document_chars:
            truncated = True
            break
    return _rows_to_markdown(rows), truncated


def _parse_xlsx(blob: bytes) -> tuple[str, bool]:
    """Every sheet as its own Markdown table.

    read_only + values_only streams rows instead of building a cell object
    graph, which matters because a spreadsheet's declared dimensions can be
    far larger than its actual data.
    """
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(blob), read_only=True, data_only=True)
    parts: list[str] = []
    acc = 0
    truncated = False
    try:
        for sheet in wb.worksheets:
            rows: list[list[str]] = []
            for values in sheet.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in values]
                if not any(c.strip() for c in cells):
                    continue  # skip the empty padding rows spreadsheets carry
                rows.append(cells)
                acc += sum(len(c) for c in cells)
                if acc >= settings.max_document_chars:
                    truncated = True
                    break
            if rows:
                parts.append(f"## {sheet.title}\n\n{_rows_to_markdown(rows)}")
            if truncated:
                break
    finally:
        wb.close()
    return "\n\n".join(parts), truncated


def _parse_pptx(blob: bytes) -> tuple[str, int, bool]:
    """Slide text, one section per slide. Returns (text, slide_count, truncated)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(blob))
    parts: list[str] = []
    acc = 0
    truncated = False
    slides = list(prs.slides)
    for i, slide in enumerate(slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            # Not every shape holds text (images, connectors, tables).
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if text:
                lines.append(text)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                lines.append(f"_Speaker notes:_ {notes}")
        if lines:
            body = "\n\n".join(lines)
            parts.append(f"## Slide {i}\n\n{body}")
            acc += len(body)
            if acc >= settings.max_document_chars:
                truncated = True
                break
    return "\n\n".join(parts), len(slides), truncated


def _parse_epub(blob: bytes) -> tuple[str, str, bool]:
    """EPUB to Markdown. Returns (title, text, truncated)."""
    import ebooklib
    from ebooklib import epub

    book = epub.read_epub(io.BytesIO(blob))
    title = ""
    meta = book.get_metadata("DC", "title")
    if meta:
        title = str(meta[0][0])

    parts: list[str] = []
    acc = 0
    truncated = False
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        html = item.get_content().decode("utf-8", errors="replace")
        md = html_to_md(html, heading_style="ATX", bullets="-").strip()
        if not md:
            continue
        parts.append(md)
        acc += len(md)
        if acc >= settings.max_document_chars:
            truncated = True
            break
    return title, "\n\n".join(parts), truncated


# Archive guards. A zip bomb is small on disk and enormous once expanded, so
# the listing never extracts anything and refuses to trust the declared sizes.
_MAX_ARCHIVE_ENTRIES = 2000


def _parse_archive(blob: bytes, source: str) -> tuple[str, bool]:
    """List an archive's contents WITHOUT extracting. Returns (text, truncated).

    Deliberately a listing, not an extraction: decompressing untrusted archive
    members is how zip bombs win. Callers who want a member can fetch it by its
    own URL, or download the archive and open it themselves.
    """
    import tarfile
    import zipfile

    rows: list[list[str]] = [["path", "size", "compressed"]]
    truncated = False
    total_declared = 0

    def add(path: str, size: int, comp: int) -> bool:
        nonlocal total_declared, truncated
        rows.append([path, _human_bytes(size), _human_bytes(comp)])
        total_declared += size
        if len(rows) > _MAX_ARCHIVE_ENTRIES:
            truncated = True
            return False
        return True

    buf = io.BytesIO(blob)
    try:
        if zipfile.is_zipfile(buf):
            buf.seek(0)
            with zipfile.ZipFile(buf) as zf:
                for info in zf.infolist():
                    if not add(info.filename, info.file_size, info.compress_size):
                        break
        else:
            buf.seek(0)
            # 'r:*' lets tarfile pick the compression; the stream is already
            # capped by max_response_bytes upstream.
            with tarfile.open(fileobj=buf, mode="r:*") as tf:
                for member in tf:
                    if not member.isfile():
                        continue
                    if not add(member.name, member.size, member.size):
                        break
    except (zipfile.BadZipFile, tarfile.TarError, EOFError) as exc:
        return f"Could not read archive: {exc}", False

    ratio = total_declared / max(len(blob), 1)
    header = (
        f"Archive listing for `{source}` — {len(rows) - 1} files, "
        f"{_human_bytes(total_declared)} uncompressed "
        f"({ratio:.0f}x compression)."
    )
    if ratio > 100:
        header += (
            "\n\n> This expands more than 100x. Treat it as untrusted and do "
            "not extract it without checking what is inside."
        )
    header += "\n\n> Contents are listed, not extracted."
    return f"{header}\n\n{_rows_to_markdown(rows)}", truncated


def _human_bytes(n: int) -> str:
    step = 1024.0
    value = float(n)
    for unit in ("B", "KB", "MB", "GB"):
        if value < step or unit == "GB":
            return f"{value:.0f} {unit}" if unit == "B" else f"{value:.1f} {unit}"
        value /= step
    return f"{value:.1f} GB"


async def _read_remote(url: str) -> tuple[bytes, str | None]:
    """SSRF-guarded download of a remote document body.

    The redirect/caps loop lives in httpfetch.httpx_stream_capped; the client
    is constructed HERE (after the guard, before any socket) so tests can keep
    monkeypatching ``httpx.AsyncClient``.
    """
    await assert_url_allowed_async(url)
    async with httpx.AsyncClient(**httpx_client_kwargs()) as client:
        _, ctype, body = await httpx_stream_capped(client, url, raise_for_status=True)
        return body, ctype or None


def _slice(full: str, start: int, length: int | None) -> tuple[str, bool, int, int]:
    """Slice [start:start+length] then smart-truncate to settings.max_content_chars.

    Returns ``(sliced_content, truncated, returned_chars, clamped_start)``.

    Invariants the caller depends on:
      * ``returned_chars`` counts SOURCE characters consumed — NOT len(content).
        smart_truncate may append a "[…truncated]" marker that is absent from
        the source, so paginating by ``start + returned_chars`` lands exactly
        on the next un-read source character (no gap, no overlap).
      * ``clamped_start`` is ``start`` clamped into ``[0, len(full)]`` so a
        caller passing a start past EOF sees where the read actually began.
      * ``truncated`` is True only when content was actually withheld: either a
        soft (smart_truncate) cut, or the slice ended before EOF.
    """
    if length is not None and length < 0:
        raise ValueError(f"length must be >= 0, got {length}")

    start = max(0, min(start, len(full)))
    end = len(full) if length is None else min(len(full), start + length)
    chunk = full[start:end]

    truncated_chunk, soft_trunc = smart_truncate(chunk, settings.max_content_chars)
    if soft_trunc:
        # smart_truncate cut the slice AND appended a marker. The real number of
        # SOURCE chars consumed is the pre-marker length, which equals what
        # smart_truncate kept before adding its suffix. Recover it by counting
        # how much of the original `chunk` survived: the kept-prefix length.
        consumed = _source_chars_consumed(chunk)
        returned_chars = consumed
        # End of consumed source for the "did we reach EOF?" decision.
        soft_end = start + consumed
    else:
        returned_chars = len(chunk)
        soft_end = end

    truncated = soft_trunc or soft_end < len(full)
    return truncated_chunk, truncated, returned_chars, start


def _source_chars_consumed(chunk: str) -> int:
    """How many leading SOURCE chars smart_truncate kept (marker excluded).

    Mirrors smart_truncate's boundary logic to recover the pre-marker length,
    so pagination by returned_chars never skips real characters.
    """
    max_chars = settings.max_content_chars
    if len(chunk) <= max_chars:
        return len(chunk)
    head = chunk[:max_chars]
    floor = int(max_chars * 0.7)
    best = -1
    # Same boundary set/logic as formatting.smart_truncate.
    for sep in ("\n\n", "\n", "。", ". ", "！", "! ", "？", "? "):
        idx = head.rfind(sep)
        if idx >= floor and idx + len(sep) > best:
            best = idx + len(sep)
    if best <= 0:
        # Hard cut at max_chars (smart_truncate did head.rstrip() + " …").
        return len(head.rstrip())
    return len(head[:best].rstrip())


def _resolve_local_path(source: str) -> Path:
    """Resolve a local-file source under the opt-in sandbox, or refuse.

    Sandbox policy (chosen by the user):
      * ``file://`` scheme is rejected outright.
      * If ``settings.document_root`` is None, local reads are DISABLED — the
        operator opts in by pointing SEARCH_MCP_DOCUMENT_ROOT at a directory.
      * Otherwise the path is resolved and must stay inside document_root;
        traversal/escape (``../../etc/passwd``, absolute paths outside the
        root, symlink escapes after resolve()) raise.
    """
    parsed = urlparse(source)
    if parsed.scheme == "file":
        raise ValueError(
            "file:// URLs are not allowed for local reads. Pass a plain path "
            "inside SEARCH_MCP_DOCUMENT_ROOT instead."
        )

    root = settings.document_root
    if root is None:
        raise PermissionError(
            "Local file reads are disabled; set SEARCH_MCP_DOCUMENT_ROOT to a "
            "directory to opt in. Remote http(s) sources are unaffected."
        )

    root = Path(root).expanduser().resolve()
    candidate = Path(source).expanduser()
    if not candidate.is_absolute():
        # Relative paths resolve against the sandbox root, not the CWD.
        candidate = root / candidate
    candidate = candidate.resolve()

    if not candidate.is_relative_to(root):
        raise PermissionError(
            f"Refusing to read {source!r}: resolves outside the document_root "
            f"sandbox ({root})."
        )
    if not candidate.exists():
        raise FileNotFoundError(source)
    return candidate


async def read_document(
    source: str,
    *,
    start: int = 0,
    length: int | None = None,
) -> DocumentResult:
    parsed = urlparse(source)
    if parsed.scheme in ("http", "https"):
        blob, ctype = await _read_remote(source)
        fmt = _detect_format(source, ctype)
    else:
        path = _resolve_local_path(source)
        blob = path.read_bytes()
        fmt = _detect_format(str(path))

    title = ""
    pages: int | None = None
    doc_truncated = False
    if fmt in ("pdf", "docx"):
        # The .pdf/.docx extension (or content-type) said binary, but the body
        # may actually be HTML — a login wall, soft-404, paywall, or CDN error
        # page served with the document's URL. Parsing those bytes as PDF/DOCX
        # raises; degrade to HTML extraction instead of crashing so the caller
        # still sees the (useful) wall/error text.
        try:
            if fmt == "pdf":
                title, full, pages, doc_truncated = _parse_pdf(blob)
            else:
                full, doc_truncated = _parse_docx(blob)
        except Exception as e:
            log.info("%s parse failed for %s; falling back to HTML: %s", fmt, source, e)
            fmt = "html"
            pages = None
            full = _parse_html(blob)
    elif fmt in ("xlsx", "pptx", "epub"):
        # Same soft-failure story as pdf/docx: the extension promised an Office
        # container but the body may be an HTML wall or error page.
        try:
            if fmt == "xlsx":
                full, doc_truncated = _parse_xlsx(blob)
            elif fmt == "pptx":
                full, pages, doc_truncated = _parse_pptx(blob)
            else:
                title, full, doc_truncated = _parse_epub(blob)
        except Exception as e:
            log.info("%s parse failed for %s; falling back to HTML: %s", fmt, source, e)
            fmt = "html"
            pages = None
            full = _parse_html(blob)
    elif fmt == "csv":
        full, doc_truncated = _parse_csv(blob)
    elif fmt == "archive":
        full, doc_truncated = _parse_archive(blob, source)
    elif fmt == "code":
        full = _parse_code(blob, source)
    elif fmt == "html":
        full = _parse_html(blob)
    elif fmt in ("text", "markdown"):
        full = _parse_text(blob)
    elif fmt == "image":
        # read_doc extracts text; an image has none. Point at the tool that
        # can actually describe it rather than returning an empty document.
        raise ValueError(
            f"{source!r} is an image. Use `fetch` to get its dimensions, format "
            "and size, or `fetch(inline=True)` to pass the image itself to a "
            "vision-capable model."
        )
    else:
        raise ValueError(
            f"Unsupported document format for {source!r}. Supported: pdf, docx, "
            "xlsx, pptx, epub, csv, archive (zip/tar), code, html, text, markdown."
        )

    content, slice_truncated, returned, clamped_start = _slice(full, start, length)
    # truncated if EITHER the parse hit a bomb-cap OR this slice withheld text.
    return DocumentResult(
        source=source,
        format=fmt,
        title=title,
        content=content,
        truncated=slice_truncated or doc_truncated,
        pages=pages,
        tokens_estimated=estimate_tokens(content),
        total_chars=len(full),
        start=clamped_start,
        returned_chars=returned,
    )
